"""
routes/upload.py — Document Upload & Learning Package Generation
=================================================================
Handles file uploads (PDF, DOCX, TXT, PPTX) and text-paste submissions.
Orchestrates a two-model AI pipeline on one Ollama host:
  - gpt-oss   → summary generation (async via httpx or run_in_executor)
  - deepseek  → quiz generation (background thread, polled via quiz-status)
Also includes PDF quality assessment, text extraction helpers, and a
deterministic fallback payload when AI services are unavailable.

Cross-reference: routes/content.py (results/modules), routes/history.py (persistence).
"""

import io
import os
import zipfile
from datetime import datetime, timezone
from typing import Optional

from fastapi import APIRouter, Depends, File, Form, Request, UploadFile

from backend.utils.api_errors import message_error
from backend.utils.rate_limit import limiter
from backend.utils.sanitization import sanitize_multiline_text, sanitize_single_line
from backend.database.db import history_collection
from backend.middleware.auth_middleware import get_current_user
from backend.services.ollama_service import generate_learning_package
from backend.services.ai_service import AIService
from backend.services.schemas import SummaryRequest, QuizRequest, SummaryResponse

# Shared AI service instance with separate summary and quiz model clients
_ai_service = AIService()


async def _async_summarize(doc_title: str, extracted_text: str):
    """
    Call the summary model using httpx async — avoids run_in_executor
    threading conflicts with Ollama's OLLAMA_NUM_PARALLEL=1 queue.
    Returns a SummaryResponse or raises on failure.
    """
    import httpx
    import asyncio
    import json as _json
    import os
    from backend.services.schemas import SummaryResponse
    from backend.services.ollama_client import SummaryOllamaSettings

    settings = SummaryOllamaSettings()
    # Use more text for PPTX (slide text is dense and short per slide)
    # 8000 chars covers ~20 slides; regular docs get 4000 chars
    max_chars = 8000 if "pptx" in doc_title.lower() or len(extracted_text) > 6000 else 4000
    text = extracted_text[:max_chars].strip()

    prompt = f"""OUTPUT ONLY VALID JSON. Start with {{ immediately.

{{"summary_title":"<title max 80 chars>","authors":"<names or Unknown authors>","overview":"<2-3 sentences on purpose and main argument>","body":["<paragraph 1: background 60-100 words>","<paragraph 2: methods/findings 60-100 words>","<paragraph 3: implications 60-100 words>"],"takeaways":["<specific claim 1>","<specific claim 2>","<specific claim 3>","<specific claim 4>","<specific claim 5>"],"topics":["topic1","topic2","topic3","topic4"],"chunks_used":1}}

Fill in all fields based on this document.
Title: {doc_title}
Content: {text}"""

    # ── Bridge mode: route through EC2 bridge instead of direct Tailscale ──
    bridge_url = os.getenv("BRIDGE_URL", "").strip()
    if bridge_url:
        async with httpx.AsyncClient(timeout=300) as client:
            resp = await client.post(f"{bridge_url}/summary", json={"prompt": prompt})
            resp.raise_for_status()
            raw = resp.json()
            response_text = raw.get("response", "")
            if not response_text:
                raise ValueError("Bridge returned empty response")
            parsed = _json.loads(response_text)
            return SummaryResponse.model_validate(parsed)

    payload = {
        "model": settings.model,
        "prompt": prompt,
        "stream": False,
        "format": "json",
        "think": False,
        "options": {
            "temperature": 0.2,
            "num_predict": 4096,  # gpt-oss ignores think:False, needs room for thinking + response
            "num_ctx": 8192,      # ensure enough context window
        },
    }

    last_error = None
    for attempt in range(4):
        try:
            wait = [0, 3, 6, 10][attempt]
            if wait:
                await asyncio.sleep(wait)
            async with httpx.AsyncClient(timeout=max(settings.timeout_seconds, 120)) as client:
                resp = await client.post(f"{settings.base_url}/api/generate", json=payload)
                resp.raise_for_status()
                raw = resp.json()
            # Check for Ollama-level error
            if raw.get("error"):
                last_error = ValueError(f"Ollama error: {raw['error']}")
                print(f"[summary] attempt {attempt+1} ollama error: {raw['error']} — retrying")
                continue
            response_text = raw.get("response", "").strip()
            print(f"[summary] attempt {attempt+1} raw keys={list(raw.keys())} response_len={len(response_text)} done={raw.get('done')} done_reason={raw.get('done_reason')}")
            if response_text:
                break
            print(f"[summary] attempt {attempt+1} empty — retrying")
            last_error = ValueError("Summary model returned an empty response")
        except Exception as e:
            print(f"[summary] attempt {attempt+1} error: {repr(e)} — retrying")
            last_error = e
    else:
        raise last_error or ValueError("Summary model failed after 4 attempts")

    # Clean markdown fences
    import re as _re
    cleaned = response_text
    if cleaned.startswith("```"):
        cleaned = _re.sub(r"^```[a-zA-Z]*\n?", "", cleaned)
        cleaned = _re.sub(r"\n?```\s*$", "", cleaned).strip()

    # Extract JSON object
    start = cleaned.find("{")
    end = cleaned.rfind("}")
    if start != -1 and end != -1:
        cleaned = cleaned[start:end + 1]

    # Repair truncated JSON
    try:
        data = _json.loads(cleaned)
    except _json.JSONDecodeError:
        import re as _re2
        s = cleaned.rstrip()
        in_string = False
        i = 0
        while i < len(s):
            c = s[i]
            if c == '\\' and in_string:
                i += 2
                continue
            if c == '"':
                in_string = not in_string
            i += 1
        if in_string:
            s += '"'
        s = _re2.sub(r',\s*$', '', s.rstrip())
        s += ']' * max(s.count('[') - s.count(']'), 0)
        s += '}' * max(s.count('{') - s.count('}'), 0)
        data = _json.loads(s)
        data = _json.loads(cleaned)

    # Normalize authors to string
    if isinstance(data.get("authors"), list):
        data["authors"] = ", ".join(str(a) for a in data["authors"]) or "Unknown authors"

    # Ensure required fields have safe defaults
    data.setdefault("summary_title", doc_title)
    data.setdefault("authors", "Unknown authors")
    data.setdefault("overview", "")
    data.setdefault("body", [])
    data.setdefault("takeaways", [])
    data.setdefault("topics", [])
    data.setdefault("chunks_used", 1)

    return SummaryResponse.model_validate(data)

# ── Background quiz job store ─────────────────────────────────────────────────
# In-memory dict that maps job_id -> quiz status/result.
# Populated by _run_quiz_in_background, polled by the /upload/quiz-status endpoint.
# NOTE: Not persisted across server restarts — quiz generation is best-effort.
import threading as _threading
_upload_quiz_jobs: dict = {}  # { job_id: {"status":"pending"|"done"|"error", "quiz":[], "error":str} }


def _run_quiz_in_background(job_id: str, doc_title: str, summary_response, history_id: str) -> None:
    """
    Thread target — generates quiz on Mac 2 (deepseek), stores result in the
    in-memory job dict, then persists to MongoDB via a sync pymongo client
    (since this runs in a daemon thread without access to the async event loop).

    Cross-reference: The upload endpoint fires this via `threading.Thread`.
                     The frontend polls /upload/quiz-status/{job_id}.
    """
    import time as _t, asyncio as _asyncio
    try:
        t0 = _t.time()
        model = _ai_service.quiz_client.settings.model
        print(f"[quiz] Quiz model generating — model={model} doc={doc_title}")
        quiz_response = _ai_service.generate_quiz(
            QuizRequest(title=doc_title, summary=summary_response,
                        question_count=8, difficulty="intermediate")
        )
        elapsed = round(_t.time() - t0, 1)
        model = _ai_service.quiz_client.settings.model
        # ── Convert response to legacy frontend format ──
        quiz_data = []
        for q in quiz_response.questions:
            try:
                quiz_data.append({
                    "q": q.question if hasattr(q, "question") else str(q),
                    "opts": q.options if hasattr(q, "options") else [],
                    "correct": q.correct_index if hasattr(q, "correct_index") else 0,
                    "explanation": q.explanation if hasattr(q, "explanation") else "",
                    "topic": q.topic if hasattr(q, "topic") else "General",
                })
            except Exception as _qe:
                print(f"[quiz] skipping malformed question: {_qe}")
                continue
        print(f"[upload] Quiz model done — model={model} time={elapsed}s questions={len(quiz_data)} doc={doc_title}")
        _upload_quiz_jobs[job_id] = {"status": "done", "quiz": quiz_data, "error": None}

        # ── Persist quiz to MongoDB using pymongo sync client (thread-safe) ──
        try:
            from bson import ObjectId
            import os
            from pymongo import MongoClient
            mongo_url = os.getenv("MONGO_URL") or os.getenv("MONGODB_URL", "mongodb://localhost:27017")
            db_name   = os.getenv("DATABASE_NAME") or os.getenv("MONGODB_DB_NAME", "learnova")
            sync_client = MongoClient(mongo_url, serverSelectionTimeoutMS=10000)
            sync_db = sync_client[db_name]
            sync_db["history"].update_one(
                {"_id": ObjectId(history_id)},
                {"$set": {"quizFull": quiz_data, "quizData": quiz_data, "total": len(quiz_data)}}
            )
            sync_client.close()
            print(f"[upload] ✅ Quiz persisted to history: {history_id}")
        except Exception as save_err:
            print(f"[upload] ⚠️  Quiz generated but MongoDB save failed: {save_err}")
    except Exception as exc:
        print(f"[upload] Quiz model failed: {repr(exc)}")
        _upload_quiz_jobs[job_id] = {"status": "error", "quiz": [], "error": str(exc)}

# ── In-memory quiz job store (declaration repeated for module-level access) ───
# { job_id: { "status": "pending"|"done"|"error", "quiz": [...], "error": str } }
_upload_quiz_jobs: dict = {}

router = APIRouter()

# ── Upload constraints ────────────────────────────────────────────────────────
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB

# Free-tier users are restricted to text-based formats; PPTX requires Pro.
ALLOWED_EXTENSIONS_FREE = {"pdf", "txt", "doc", "docx"}
ALLOWED_EXTENSIONS_PRO = {"pdf", "txt", "doc", "docx", "ppt", "pptx"}
_OLE_SIGNATURE = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


# ----------------------------- Text Extraction Helpers -----------------------------
def _extract_text_from_pdf(data: bytes) -> str:
    """Extract readable text from PDF bytes."""
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(data))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages)
    except Exception as e:
        print(f"[upload] PDF extraction error: {e}")
        return ""


# ── PDF quality checker ──────────────────────────────────────────────────────
# Thresholds — tune these if needed
_MIN_CHARS_PER_PAGE     = 80    # fewer than this per page = likely scanned/image PDF
_MIN_WORD_LENGTH_AVG    = 3.0   # avg word length below this = likely OCR garbage
_MAX_GARBAGE_RATIO      = 0.35  # more than 35% non-ASCII/non-printable = corrupted
_MIN_READABLE_PAGES_PCT = 0.40  # at least 40% of pages must pass quality check


def _assess_pdf_quality(data: bytes, extracted_text: str, page_count: int) -> tuple[bool, str]:
    """
    Returns (is_readable, reason).
    is_readable=False means the PDF is scanned, blurry, or image-only.
    """
    if page_count == 0:
        return False, "empty"

    total_chars = len(extracted_text.strip())
    chars_per_page = total_chars / max(page_count, 1)

    # Check 1 — almost no text extracted at all
    if total_chars < 50:
        return False, "no_text"

    if chars_per_page < _MIN_CHARS_PER_PAGE:
        return False, "too_sparse"

    # Check 2 — garbage character ratio (non-printable, replacement chars)
    import unicodedata
    garbage = sum(
        1 for ch in extracted_text
        if (not ch.isprintable() and ch not in ("\n", "\t", "\r"))
        or ch == "�"   # Unicode replacement character — sign of bad encoding
        or (ord(ch) > 127 and unicodedata.category(ch) == "So")  # misc symbols
    )
    garbage_ratio = garbage / max(total_chars, 1)
    if garbage_ratio > _MAX_GARBAGE_RATIO:
        return False, "garbage_text"

    # Check 3 — word length average (real text has avg ~4-7 chars per word)
    words = [w for w in extracted_text.split() if w.isalpha()]
    if words:
        avg_len = sum(len(w) for w in words) / len(words)
        if avg_len < _MIN_WORD_LENGTH_AVG and len(words) > 20:
            return False, "garbled_words"

    # Check 4 — per-page extraction (need at least 40% of pages to have real text)
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(data))
        readable_pages = sum(
            1 for page in reader.pages
            if len((page.extract_text() or "").strip()) >= _MIN_CHARS_PER_PAGE
        )
        readable_pct = readable_pages / max(len(reader.pages), 1)
        if readable_pct < _MIN_READABLE_PAGES_PCT:
            return False, "mostly_images"
    except Exception:
        pass  # If PyPDF2 fails here, we already have text so it's probably OK

    return True, "ok"


def _build_blur_error(reason: str) -> dict:
    """Return a structured error payload the frontend can display."""
    messages = {
        "no_text":       "No readable text found in this PDF.",
        "too_sparse":    "This PDF appears to be a scanned document — very little text could be extracted.",
        "garbage_text":  "This PDF contains mostly unreadable characters, likely from a bad scan or image conversion.",
        "garbled_words": "The text extracted from this PDF appears garbled — it may be a low-quality scan.",
        "mostly_images": "Most pages in this PDF are images with no selectable text.",
        "empty":         "This PDF appears to be empty.",
    }
    detail = messages.get(reason, "This PDF could not be read properly.")
    return {
        "error": "unreadable_pdf",
        "detail": detail,
        "user_message": (
            f"{detail} "
            "Please upload a clearer version, a text-based PDF, or paste the text directly using the text input."
        ),
        "suggestions": [
            "Use a text-based PDF (not a scan)",
            "Copy and paste the text directly",
            "Try converting the PDF to a Word document first",
            "If it is a scanned document, run OCR software on it first",
        ],
    }


def _extract_text_from_docx(data: bytes) -> str:
    """Extract readable text from DOCX bytes."""
    try:
        import docx
        doc = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        print(f"[upload] DOCX extraction error: {e}")
        return ""


def _extract_text_from_pptx(data: bytes) -> str:
    """Extract readable text from PPT/PPTX bytes."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        print(f"[upload] PPTX extraction error: {e}")
        return ""


def _extract_text(filename: str, data: bytes) -> tuple[str, str, int]:
    """Return extracted text, normalized file type label, and estimated page count."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "txt"

    if ext == "pdf":
        text = _extract_text_from_pdf(data)
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(io.BytesIO(data))
            pages = len(reader.pages)
        except Exception:
            pages = max(1, len(text) // 3000)
        return text, "PDF", max(1, pages)

    elif ext in ("doc", "docx"):
        text = _extract_text_from_docx(data)
        pages = max(1, len(text) // 3000)
        return text, "DOCX", pages

    elif ext in ("ppt", "pptx"):
        text = _extract_text_from_pptx(data)
        # Count actual slides for accurate page count
        try:
            from pptx import Presentation
            import io as _io
            prs = Presentation(_io.BytesIO(data))
            pages = max(1, len(prs.slides))
        except Exception:
            pages = max(1, len(text) // 500)
        return text, "PPTX", pages

    else:
        # Plain text / fallback
        try:
            text = data.decode("utf-8", errors="replace")
        except Exception:
            text = ""
        pages = max(1, len(text) // 3000)
        return text, ext.upper(), pages


def _zip_contains_prefix(data: bytes, prefix: str) -> bool:
    """Return True when a ZIP-based Office file contains an expected folder."""
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            return any(name.startswith(prefix) for name in archive.namelist())
    except zipfile.BadZipFile:
        return False


def _validate_uploaded_file(filename: str, data: bytes) -> str | None:
    """Validate uploaded file bytes against the extension before extraction."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    head = data[:16]

    if ext == "pdf":
        if not data.lstrip().startswith(b"%PDF"):
            return "Uploaded .pdf file is not a valid PDF document."
    elif ext == "txt":
        if b"\x00" in data[:4096]:
            return "Uploaded .txt file appears to contain binary data."
        try:
            data[:4096].decode("utf-8")
        except UnicodeDecodeError:
            return "Uploaded .txt file must be valid UTF-8 text."
    elif ext == "docx":
        if not _zip_contains_prefix(data, "word/"):
            return "Uploaded .docx file is not a valid Word document."
    elif ext == "pptx":
        if not _zip_contains_prefix(data, "ppt/"):
            return "Uploaded .pptx file is not a valid PowerPoint document."
    elif ext == "doc":
        if not (head.startswith(_OLE_SIGNATURE) or _zip_contains_prefix(data, "word/")):
            return "Uploaded .doc file is not a valid Word document."
    elif ext == "ppt":
        if not (head.startswith(_OLE_SIGNATURE) or _zip_contains_prefix(data, "ppt/")):
            return "Uploaded .ppt file is not a valid PowerPoint document."

    return None


async def _get_next_seq_id(user_id: str) -> int:
    """Generate the next per-user sequential ID for history records."""
    count = await history_collection.count_documents({"userId": user_id})
    return count + 1


def _build_fallback_payload(doc_title: str, page_count: int) -> dict:
    """Build deterministic fallback content when AI service is unavailable."""
    return {
        "summary": {
            "body": [
                f"This document titled '{doc_title}' has been processed by Learnova's AI engine. It spans {page_count} {'page' if page_count == 1 else 'pages'} covering key academic concepts.",
                "The themes identified include core methodology, evidence base, and author conclusions — forming the basis for the comprehension quiz below.",
            ],
            "takeaways": [
                "The document presents a structured argument supported by academic references.",
                "Core findings are framed within a broader context with practical implications.",
                "Review this summary carefully — the quiz will test understanding of the main claims.",
            ],
        },
        "quiz": [
            {"q": "What is the primary focus of this document?", "opts": ["Empirical measurement", "Theoretical critique", "Systematic synthesis", "Policy evaluation"], "correct": 2, "explanation": "The document takes a synthesis approach."},
            {"q": "Which methodology is primarily employed?", "opts": ["Randomised trial", "Ethnographic fieldwork", "Literature review", "Longitudinal survey"], "correct": 2, "explanation": "A literature review and analytical framework are central."},
            {"q": "What type of evidence is most heavily used?", "opts": ["Anecdotal reports", "Peer-reviewed studies", "Government statistics", "Industry benchmarks"], "correct": 1, "explanation": "Peer-reviewed studies form the backbone of evidence."},
            {"q": "Who is the target audience?", "opts": ["General public", "Academic researchers", "Policy makers only", "Undergrad students"], "correct": 1, "explanation": "Technical language indicates researchers and practitioners."},
            {"q": "What gap is identified in existing work?", "opts": ["Lack of data", "Under-representation", "Insufficient longitudinal research", "Overemphasis on theory"], "correct": 2, "explanation": "Longitudinal evidence is underdeveloped in this field."},
            {"q": "What does the document recommend?", "opts": ["Abandon frameworks", "Cross-disciplinary collaboration", "Quantitative only", "Single institution studies"], "correct": 1, "explanation": "Cross-disciplinary collaboration is most promising."},
            {"q": "Which factor most influences outcomes?", "opts": ["Funding levels", "Institutional support", "Individual motivation", "Technology availability"], "correct": 1, "explanation": "Institutional support is the dominant conditioning factor."},
            {"q": "What is the overall contribution?", "opts": ["Definitive theory proof", "Synthesised framework", "Replication study", "Full critique"], "correct": 1, "explanation": "A synthesised framework organises complex findings."},
        ],
        "analysis": {
            "strengths": ["Core understanding", "Concept linkage", "Evidence awareness"],
            "weaknesses": ["Application depth", "Long-term retention"],
            "recommendations": ["Review core concepts regularly", "Practice applied questions", "Revisit evidence summaries"],
            "studyNext": ["Related academic literature", "Applied case studies", "Cross-disciplinary research"],
        },
        "modules": [
            {"title": "Introduction to the Topic", "type": "youtube", "url": f"https://www.youtube.com/results?search_query={doc_title.replace(' ', '+')}", "description": "Video overview of the main topic"},
            {"title": "Academic Overview", "type": "google", "url": f"https://www.google.com/search?q={doc_title.replace(' ', '+')}+academic", "description": "Academic resources on this topic"},
            {"title": "Key Concepts Explained", "type": "youtube", "url": f"https://www.youtube.com/results?search_query={doc_title.replace(' ', '+')}+explained", "description": "Concept explanations"},
            {"title": "Further Reading", "type": "google", "url": f"https://www.google.com/search?q={doc_title.replace(' ', '+')}+research", "description": "Deeper research materials"},
            {"title": "Related Topics", "type": "youtube", "url": f"https://www.youtube.com/results?search_query={doc_title.replace(' ', '+')}+tutorial", "description": "Tutorials on related topics"},
        ],
    }


async def _generate_ai_payload(
    doc_title: str,
    file_type: str,
    extracted_text: str,
    fallback_payload: dict,
) -> dict:
    """
    Generate AI payload using dual-Mac routing:
    - Mac 1 (gpt-oss)     → summary via AIService
    - Mac 2 (deepseek)    → quiz via AIService quiz_client
    Falls back to legacy single-Mac path if both fail.
    """
    import asyncio

    # ── Step 1: Summary on Mac 1 (gpt-oss) ──────────────────────────────────
    summary_response = None
    try:
        loop = asyncio.get_event_loop()
        import time as _time
        _t0 = _time.time()
        summary_response = await loop.run_in_executor(
            None,
            lambda: _ai_service.summarize(
                SummaryRequest(title=doc_title, text=extracted_text)
            )
        )
        _summary_time = round(_time.time() - _t0, 1)
        summary_model = _ai_service.client.settings.model
        print(f"[upload] Summary model done — model={summary_model} time={_summary_time}s doc={doc_title}")
    except Exception as e:
        print(f"[upload] Summary model failed: {repr(e)} — trying legacy path")

    # ── Step 2: Quiz on Mac 2 (deepseek) ────────────────────────────────────
    quiz_data = None
    if summary_response:
        try:
            loop = asyncio.get_event_loop()
            _tq = _time.time()
            quiz_response = await loop.run_in_executor(
                None,
                lambda: _ai_service.generate_quiz(
                    QuizRequest(
                        title=doc_title,
                        summary=summary_response,
                        question_count=8,
                        difficulty="intermediate",
                    )
                )
            )
            _quiz_time = round(_time.time() - _tq, 1)
            quiz_model = _ai_service.quiz_client.settings.model
            # Convert to legacy quiz format expected by frontend
            quiz_data = [
                {
                    "q": q.question,
                    "opts": q.options,
                    "correct": q.correct_index,
                    "explanation": q.explanation,
                    "topic": q.topic,
                }
                for q in quiz_response.questions
            ]
            print(f"[upload] Quiz model done — model={quiz_model} time={_quiz_time}s questions={len(quiz_data)} doc={doc_title}")
        except Exception as e:
            print(f"[upload] Quiz model failed: {repr(e)}")

    # ── Step 3: Build payload if both succeeded ──────────────────────────────
    if summary_response and quiz_data:
        return {
            "summary": {
                "body": summary_response.body,
                "takeaways": summary_response.takeaways,
                "title": summary_response.summary_title,
                "authors": summary_response.authors,
                "topics": summary_response.topics,
            },
            "analysis": {
                "strengths": [],
                "weaknesses": [],
                "studyNext": [],
            },
            "modules": [],
            "quiz": quiz_data,
        }

    # ── Step 4: Legacy single-Mac fallback ───────────────────────────────────
    print(f"[upload] Falling back to legacy path for: {doc_title}")
    try:
        payload = await generate_learning_package(
            title=doc_title,
            file_type=file_type,
            text_content=extracted_text,
        )
        print(f"[upload] Legacy path succeeded for: {doc_title}")
        return payload
    except Exception as error:
        print(f"[upload] All AI paths failed, using fallback: {repr(error)}")
        return fallback_payload


def _build_history_doc(
    user_id: str,
    seq_id: int,
    doc_title: str,
    file_type: str,
    page_count: int,
    word_estimate: int,
    ai_payload: dict,
    now: datetime,
) -> dict:
    """Build the MongoDB history document from normalized upload data."""
    return {
        "userId": user_id,
        "seqId": seq_id,
        "title": doc_title,
        "fileType": file_type,
        "pageCount": page_count,
        "wordEstimate": word_estimate,
        "authors": ai_payload["summary"].get("authors", "Unknown authors"),
        "summary": ai_payload["summary"],
        "analysis": ai_payload["analysis"],
        "modules": ai_payload["modules"],
        "quizFull": ai_payload["quiz"],
        "done": False,
        "score": None,
        "correct": None,
        "total": 8,
        "userAnswers": [],
        "uploadedAt": now,
        "completedAt": None,
    }


def _format_processed_time(now: datetime) -> str:
    """Format processing time in the same style used by existing frontend."""
    if os.name != "nt":
        return now.strftime("%-I:%M %p").lower()
    return now.strftime("%I:%M %p").lower().lstrip("0")


def _build_upload_response(
    history_id: str,
    seq_id: int,
    doc_title: str,
    file_type: str,
    page_count: int,
    word_estimate: int,
    ai_payload: dict,
    processed_at: str,
) -> dict:
    """Build API response payload returned to upload page after processing."""
    return {
        "historyId": history_id,
        "seqId": seq_id,
        "title": doc_title,
        "fileType": file_type,
        "pageCount": page_count,
        "meta": f".{file_type.lower()} · {page_count} {'page' if page_count == 1 else 'pages'}",
        "infoRows": [
            ["File type", file_type],
            ["Estimated words", f"{word_estimate:,}"],
            ["Pages", str(page_count)],
            ["Language", "English"],
            ["Processed", processed_at],
        ],
        "summary": {
            "title": doc_title,
            "authors": "Uploaded document",
            "pages": f"{page_count} page{'s' if page_count != 1 else ''}",
            "body": ai_payload["summary"]["body"],
            "takeaways": ai_payload["summary"]["takeaways"],
            "strengths": ai_payload["analysis"]["strengths"],
            "weaknesses": ai_payload["analysis"]["weaknesses"],
            "studyNext": ai_payload["analysis"]["studyNext"],
        },
        "quizData": ai_payload["quiz"],
        "modules": ai_payload["modules"],
    }


# ----------------------------- Quiz Status Endpoint -----------------------------

# GET /upload/quiz-status/{job_id} — poll for Mac 2 quiz readiness
@router.get("/upload/quiz-status/{job_id}")
async def upload_quiz_status(job_id: str):
    """
    GET /upload/quiz-status/{job_id} — polling endpoint for background quiz generation.
    Returns "pending" while Mac 2 is still working, "done" with quiz data when ready,
    or "error" if generation failed. The frontend polls this after the initial upload
    response returns a quizJobId (see the /upload endpoint below).
    """
    job = _upload_quiz_jobs.get(job_id)
    if not job:
        return {"job_id": job_id, "status": "pending", "quiz": [], "error": None}
    return {"job_id": job_id, "status": job["status"], "quiz": job.get("quiz", []), "error": job.get("error")}


# ----------------------------- Upload Endpoint -----------------------------

# POST /upload — main upload handler
@router.post("/upload")
@limiter.limit("10/hour")
async def upload_document(
    request: Request,
    current_user: dict = Depends(get_current_user),
    file: Optional[UploadFile] = File(None),
    text_content: Optional[str] = Form(None),
    title: Optional[str] = Form(None),
):
    """
    POST /upload — process a file or text-paste submission.
    Flow:
      1. Validates file type against user tier (free vs pro).
      2. PDFs go through a quality gate (scanned/image PDF detection).
      3. Extracts text via format-specific helpers.
      4. Phase 1 — Mac 1 generates summary (via _async_summarize).
      5. Phase 2 — builds AI payload; falls back to legacy path.
      6. Saves history document to MongoDB immediately.
      7. Phase 4 — fires quiz generation on Mac 2 in a background thread.
      8. Returns immediately with a quizJobId for the frontend to poll.
    Cross-reference: Poll quiz progress at /upload/quiz-status/{quizJobId}.
                     Results viewed through routes/content.py.
    """
    user_id = str(current_user["_id"])
    user_tier = current_user.get("tier", "free")

    # ── Determine source: file upload OR text paste ──────────────────────────
    if file and file.filename:
        ext = file.filename.rsplit(".", 1)[-1].lower() if "." in file.filename else ""
        allowed = ALLOWED_EXTENSIONS_PRO if user_tier == "pro" else ALLOWED_EXTENSIONS_FREE

        if ext not in allowed:
            if user_tier != "pro" and ext in ("ppt", "pptx"):
                return message_error(403, "PowerPoint upload requires a Pro account. Please upgrade to continue.")
            return message_error(
                400,
                f"File type .{ext} is not supported. Allowed: {', '.join(sorted(allowed))}",
            )

        data = await file.read()
        if len(data) > MAX_FILE_SIZE:
            return message_error(413, "File exceeds 10 MB limit. Please upload a smaller file.")
        validation_error = _validate_uploaded_file(file.filename, data)
        if validation_error:
            return message_error(400, validation_error)

        raw_title = title or file.filename.rsplit(".", 1)[0]
        doc_title = sanitize_single_line(raw_title, max_length=160) or "Uploaded document"
        extracted_text, file_type, page_count = _extract_text(file.filename, data)

        # ── PDF quality gate ────────────────────────────────────────────────
        if ext == "pdf":
            is_readable, reason = _assess_pdf_quality(data, extracted_text, page_count)
            if not is_readable:
                from fastapi.responses import JSONResponse
                error_payload = _build_blur_error(reason)
                return JSONResponse(status_code=422, content=error_payload)
        # ────────────────────────────────────────────────────────────────────

    elif text_content and text_content.strip():
        extracted_text = sanitize_multiline_text(text_content)
        file_type = "TXT"
        page_count = max(1, len(extracted_text) // 3000)
        doc_title = sanitize_single_line(title, max_length=160) if title else "Pasted text"
        doc_title = doc_title or "Pasted text"

    else:
        return message_error(400, "Please upload a file or paste text content.")

    # Rough word count for display metadata; used in the frontend info rows
    word_estimate = len(extracted_text.split()) if extracted_text else page_count * 350

    import time as _t, uuid as _uuid
    fallback_payload = _build_fallback_payload(doc_title, page_count)

    # ── Phase 1: Summary on Mac 1 (gpt-oss) ─────────────────────────────────
    # Uses _async_summarize (httpx-based async call) instead of run_in_executor
    # to avoid conflicts with Ollama's OLLAMA_NUM_PARALLEL=1 queue.
    summary_response = None
    try:
        t0 = _t.time()
        model = _ai_service.client.settings.model
        print(f"[summary] Summary model running — model={model} doc={doc_title}")
        summary_response = await _async_summarize(doc_title, extracted_text)
        elapsed = round(_t.time() - t0, 1)
        print(f"[upload] Summary model done — model={model} time={elapsed}s doc={doc_title}")
    except Exception as e:
        print(f"[upload] Summary model failed: {repr(e)} — falling back to legacy")

    # ── Phase 2: Build ai_payload from summary ───────────────────────────────
    if summary_response:
        ai_payload = {
            "summary": {
                "body": summary_response.body,
                "takeaways": summary_response.takeaways,
                "title": summary_response.summary_title,
                "authors": summary_response.authors,
                "topics": summary_response.topics,
            },
            "analysis": {"strengths": [], "weaknesses": [], "studyNext": []},
            "modules": [],
            "quiz": [],
        }
    else:
        # Legacy fallback — gets summary AND quiz in one shot
        ai_payload = await _generate_ai_payload(
            doc_title=doc_title, file_type=file_type,
            extracted_text=extracted_text, fallback_payload=fallback_payload,
        )

    # ── Phase 3: Save history immediately ───────────────────────────────────
    # The history document is persisted BEFORE the quiz finishes so the frontend
    # can navigate to the results page. Quiz data is backfilled via the background
    # thread (Phase 4) or the default fallback quiz.
    now = datetime.now(timezone.utc)
    seq_id = await _get_next_seq_id(user_id)
    history_doc = _build_history_doc(
        user_id=user_id, seq_id=seq_id, doc_title=doc_title,
        file_type=file_type, page_count=page_count,
        word_estimate=word_estimate, ai_payload=ai_payload, now=now,
    )
    result = await history_collection.insert_one(history_doc)
    history_id = str(result.inserted_id)

    # ── Phase 4: Fire quiz on Mac 2 (deepseek) in background thread ──────────
    # Quiz generation is CPU-intensive and uses a separate Ollama model, so it
    # runs on a daemon thread rather than blocking the HTTP response. The frontend
    # polls /upload/quiz-status/{quiz_job_id} until the quiz is ready.
    quiz_job_id = None
    if summary_response:
        quiz_job_id = str(_uuid.uuid4())
        _upload_quiz_jobs[quiz_job_id] = {"status": "pending", "quiz": [], "error": None}
        t = _threading.Thread(
            target=_run_quiz_in_background,
            args=(quiz_job_id, doc_title, summary_response, history_id),
            daemon=True,
        )
        t.start()
        print(f"[upload] Quiz model queued — job_id={quiz_job_id}")

    # ── Phase 5: Return immediately (do NOT wait for quiz) ───────────────────
    # The response includes quizJobId for frontend polling, and quizReady=false
    # when background quiz generation was queued.
    processed_at = _format_processed_time(now)
    response = _build_upload_response(
        history_id=history_id, seq_id=seq_id, doc_title=doc_title,
        file_type=file_type, page_count=page_count,
        word_estimate=word_estimate, ai_payload=ai_payload,
        processed_at=processed_at,
    )
    response["quizJobId"] = quiz_job_id
    response["quizReady"] = quiz_job_id is None
    return response
