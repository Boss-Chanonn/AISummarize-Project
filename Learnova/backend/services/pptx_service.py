"""
pptx_service.py  —  PPTX processing pipeline
=============================================
Handles everything PPTX-specific:
  1. Per-slide text extraction (python-pptx)
  2. Per-slide AI summarisation (Mac 1 — gpt-oss)
  3. Slicing a stored summary by slide range
  4. Module unlock logic (are all slides covered?)

Pipeline flow:
  1. Upload → extract_slides() extracts text per slide
  2. summarise_all_slides() → AI summaries stored in MongoDB
  3. Quiz generation → build_range_summary_response() creates a SummaryResponse
     from a subset of slide summaries (no re-processing needed)
  4. all_slides_covered() → checks if the student has studied every slide

Install dependency if not present:
    pip install python-pptx --break-system-packages

Cross-references:
  - Uses OllamaClient with SummaryOllamaSettings (Mac 1, gpt-oss).
  - build_range_summary_response creates a SummaryResponse (see schemas.py).
"""

from __future__ import annotations

import io
import os
from dataclasses import dataclass, field

from pptx import Presentation  # pip install python-pptx

from backend.services.ollama_client import OllamaClient, OllamaError, SummaryOllamaSettings, BridgeClient, BridgeSettings
from backend.services.schemas import SummaryResponse


def _generate_slide_json(client: OllamaClient, prompt: str) -> dict:
    """
    Call Ollama with think:False so reasoning models don't waste tokens
    on chain-of-thought before writing the slide summary JSON.

    This is a specialised variant of OllamaClient.generate_json that:
      - Uses a lower temperature (0.1) for more deterministic output
      - Caps num_predict at 400 (slide summaries are short)
      - Cleans markdown code fences from the raw response
      - Falls back to _repair_truncated_json if the output is cut off

    Args:
        client: An OllamaClient instance (typically targeting Mac 1 / gpt-oss).
        prompt: The formatted prompt string for the slide.

    Returns:
        Parsed JSON dict with keys: summary_title, overview, topics, takeaways.

    Raises:
        OllamaError: If the response is empty after all retries.
        json.JSONDecodeError: If the response cannot be parsed even after repair.
    """
    import json as _json
    from backend.services.ollama_client import OllamaError, _repair_truncated_json
    import re as _re

    payload = {
        "model": client.settings.model,
        "prompt": prompt,
        "stream": False,
        "format": "json",
        "think": False,
        "options": {
            "temperature": 0.1,
            "num_predict": 400,   # slide summaries are short — 400 tokens is plenty
        },
    }
    raw_response = client._post("/api/generate", payload)
    response_text = raw_response.get("response", "").strip()
    if not response_text:
        raise OllamaError("Empty response from Ollama for slide summary")

    # Clean markdown fences that some models wrap around JSON output
    cleaned = response_text
    if cleaned.startswith("```"):
        cleaned = _re.sub(r"^```[a-zA-Z]*\n?", "", cleaned)
        cleaned = _re.sub(r"\n?```\s*$", "", cleaned).strip()
    start = cleaned.find("{")
    end = cleaned.rfind("}")
    if start != -1 and end != -1:
        cleaned = cleaned[start:end + 1]

    try:
        return _json.loads(cleaned)
    except _json.JSONDecodeError:
        # If the JSON was truncated (e.g. hit the token limit), try to repair it
        cleaned = _repair_truncated_json(cleaned)
        return _json.loads(cleaned)


# ── Per-slide data structures ─────────────────────────────────────────────────

@dataclass
class SlideText:
    """Raw text content extracted from a single PPTX slide.

    Attributes:
        slide_number: 1-indexed slide number.
        title:        The slide title text (if detected).
        body:         The remaining text content (non-title shapes).
        full_text:    Title + body concatenated (used for AI summarisation).
    """
    slide_number: int       # 1-indexed
    title: str
    body: str
    full_text: str          # title + body combined


@dataclass
class SlideSummary:
    """AI-generated summary for a single slide.

    Attributes:
        slide_number:  1-indexed slide number.
        summary_title: Short title for this slide's summary.
        overview:      2-3 sentence explanation of the slide's content.
        topics:        List of noun-phrase topics covered on this slide.
        takeaways:     Key claims or points from the slide.
    """
    slide_number: int
    summary_title: str
    overview: str
    topics: list[str] = field(default_factory=list)
    takeaways: list[str] = field(default_factory=list)


@dataclass
class PptxDocument:
    """Aggregate of all slide summaries for a single PPTX file.

    This is the top-level structure stored in MongoDB after processing.
    """
    title: str
    total_slides: int
    slide_summaries: list[SlideSummary]


# ── Text extraction ───────────────────────────────────────────────────────────

def extract_slides(data: bytes) -> list[SlideText]:
    """Extract per-slide text from PPTX binary data.

    Iterates every shape on every slide, collecting title and body text.
    The title is identified by looking for a shape named "title" (python-pptx
    convention). Falls back to using the first text shape as the title.

    Args:
        data: Raw PPTX file bytes (from file upload or GridFS).

    Returns:
        List of SlideText objects, one per slide, preserving slide order.
    """
    prs = Presentation(io.BytesIO(data))
    slides: list[SlideText] = []

    for i, slide in enumerate(prs.slides, start=1):
        title_text = ""
        body_parts: list[str] = []

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue
            # First text frame that looks like a title
            if not title_text and hasattr(shape, "name") and "title" in shape.name.lower():
                title_text = text
            else:
                body_parts.append(text)

        # Fallback: use first shape text as title
        if not title_text and body_parts:
            title_text = body_parts.pop(0)

        body = "\n".join(body_parts)
        full_text = f"{title_text}\n{body}".strip()

        slides.append(SlideText(
            slide_number=i,
            title=title_text or f"Slide {i}",
            body=body,
            full_text=full_text,
        ))

    return slides


# ── AI summarisation per slide ────────────────────────────────────────────────

_SLIDE_SUMMARY_PROMPT = """You are an expert academic summariser.

Summarise slide {slide_number} of a presentation titled "{pptx_title}".

Slide title: {slide_title}
Slide content:
{slide_text}

Return a JSON object with these exact keys:
{{
  "summary_title": "concise title for this slide (max 80 chars)",
  "overview": "2-3 sentence explanation of what this slide covers (min 60 chars)",
  "topics": ["topic1", "topic2"],
  "takeaways": ["key point 1", "key point 2", "key point 3"]
}}

Return only the JSON object. No preamble, no markdown.
"""


def _summarise_slide(
    client: OllamaClient,
    pptx_title: str,
    slide: SlideText,
    retries: int = 2,
) -> SlideSummary:
    """Call Mac 1 (gpt-oss) to summarise a single slide.

    If the slide has minimal text (< 30 chars), returns a stub summary
    without making an API call. Retries up to `retries` times on failure.

    Args:
        client:    OllamaClient targeting the summary model (Mac 1).
        pptx_title: Title of the overall presentation (for context).
        slide:     SlideText object with the slide's raw content.
        retries:   Maximum number of retry attempts after the first try.

    Returns:
        SlideSummary with AI-generated content, or a stub on failure.
    """
    import json

    if len(slide.full_text.strip()) < 30:
        # Slide is nearly empty — return a stub
        return SlideSummary(
            slide_number=slide.slide_number,
            summary_title=slide.title or f"Slide {slide.slide_number}",
            overview="This slide contains minimal text content.",
            topics=[],
            takeaways=[],
        )

    prompt = _SLIDE_SUMMARY_PROMPT.format(
        slide_number=slide.slide_number,
        pptx_title=pptx_title,
        slide_title=slide.title,
        slide_text=slide.full_text[:2000],   # cap to avoid token overflow
    )

    # Temporarily patch think:False onto the client call
    # by prepending a thin wrapper — ollama_client.generate_json reads its own payload
    errors: list[str] = []
    for _ in range(retries + 1):
        try:
            raw = _generate_slide_json(client, prompt)
            return SlideSummary(
                slide_number=slide.slide_number,
                summary_title=str(raw.get("summary_title", slide.title))[:80],
                overview=str(raw.get("overview", ""))[:600],
                topics=[str(t) for t in raw.get("topics", [])][:4],
                takeaways=[str(t) for t in raw.get("takeaways", [])][:4],
            )
        except (OllamaError, Exception) as exc:
            errors.append(str(exc))

    # All retries failed — return a stub with the error logged
    return SlideSummary(
        slide_number=slide.slide_number,
        summary_title=slide.title or f"Slide {slide.slide_number}",
        overview=f"Summary unavailable. ({errors[-1][:120]})",
        topics=[],
        takeaways=[],
    )


def summarise_all_slides(
    pptx_title: str,
    slides: list[SlideText],
    client: OllamaClient | None = None,
) -> list[SlideSummary]:
    """Summarise every slide via Mac 1 (gpt-oss) or EC2 bridge.

    Called once at PPTX upload time — the results are stored in MongoDB
    so that subsequent quiz sessions can slice cached summaries without
    re-processing.

    Args:
        pptx_title: Title of the presentation.
        slides:     List of SlideText objects from extract_slides().
        client:     Optional OllamaClient or BridgeClient; creates a default one if not provided.

    Returns:
        List of SlideSummary objects in slide order (same length as input).
    """
    if client is None:
        bridge_url = os.getenv("BRIDGE_URL", "").strip()
        client = BridgeClient(BridgeSettings(), endpoint="/summary") if bridge_url else OllamaClient(SummaryOllamaSettings())
    active_client = client
    return [_summarise_slide(active_client, pptx_title, slide) for slide in slides]
    return [_summarise_slide(active_client, pptx_title, slide) for slide in slides]


# ── Slice helpers ─────────────────────────────────────────────────────────────

def slice_summaries(
    all_summaries: list[SlideSummary],
    start: int,
    end: int,
) -> list[SlideSummary]:
    """
    Return slide summaries for a given range [start, end] (both inclusive, 1-indexed).

    This is the core of the "no re-processing" feature — instead of calling the
    AI again for every study session, we just slice the cached summaries from
    the initial upload.

    Args:
        all_summaries: Full list of SlideSummary from summarise_all_slides().
        start:         First slide number (inclusive).
        end:           Last slide number (inclusive).

    Returns:
        Filtered list of SlideSummary within the given range.
    """
    return [s for s in all_summaries if start <= s.slide_number <= end]


def build_range_summary_text(
    pptx_title: str,
    summaries: list[SlideSummary],
) -> str:
    """
    Flatten a list of SlideSummary objects into a single text block
    suitable for sending to the quiz generator.

    Args:
        pptx_title: Title of the presentation.
        summaries:  List of SlideSummary to flatten.

    Returns:
        Plain text block with slide numbers, titles, overviews, topics, and takeaways.
    """
    parts = [f"Presentation: {pptx_title}\n"]
    for s in summaries:
        parts.append(
            f"Slide {s.slide_number}: {s.summary_title}\n"
            f"{s.overview}\n"
            f"Key topics: {', '.join(s.topics)}\n"
            f"Takeaways: {'; '.join(s.takeaways)}\n"
        )
    return "\n".join(parts)


def build_range_summary_response(
    pptx_title: str,
    summaries: list[SlideSummary],
    start: int,
    end: int,
) -> SummaryResponse:
    """
    Build a SummaryResponse-compatible object from a slide range.

    Merges overviews, topics, and takeaways across all slides in the range,
    deduplicating while preserving order. This is passed to the quiz generator
    as if it were a regular document summary, allowing the quiz pipeline to
    work identically for documents and PPTX files.

    Args:
        pptx_title: Title of the presentation.
        summaries:  Slide summaries for the range (from slice_summaries()).
        start:      First slide number in the range.
        end:        Last slide number in the range.

    Returns:
        A SummaryResponse that can be used directly by QuizRequest.
    """
    all_topics: list[str] = []
    all_takeaways: list[str] = []
    overviews: list[str] = []

    for s in summaries:
        overviews.append(f"Slide {s.slide_number}: {s.overview}")
        all_topics.extend(s.topics)
        all_takeaways.extend(s.takeaways)

    # Deduplicate while preserving order
    seen: set[str] = set()
    unique_topics: list[str] = []
    for t in all_topics:
        if t not in seen:
            seen.add(t)
            unique_topics.append(t)

    seen = set()
    unique_takeaways: list[str] = []
    for t in all_takeaways:
        if t not in seen:
            seen.add(t)
            unique_takeaways.append(t)

    return SummaryResponse(
        summary_title=f"{pptx_title} — Slides {start}–{end}",
        authors="Presentation",
        overview=f"Session covering slides {start} to {end} of {pptx_title}.",
        body=overviews[:4] if overviews else [f"Slides {start}–{end} of {pptx_title}."],
        takeaways=unique_takeaways[:5] if unique_takeaways else ["Review the slides."],
        topics=unique_topics[:6] if unique_topics else ["General content"],
        chunks_used=len(summaries),
    )


def all_slides_covered(total_slides: int, sessions: list[dict]) -> bool:
    """
    Check whether every slide from 1..total_slides has been included
    in at least one completed study session.

    Used to determine if the student has finished the entire PPTX and
    should unlock the summary/quiz for the full deck.

    Args:
        total_slides: Total number of slides in the PPTX.
        sessions:     List of session dicts from MongoDB (each has
                      slide_start, slide_end, and done flag).

    Returns:
        True if every slide has been covered in a completed session.
    """
    covered: set[int] = set()
    for session in sessions:
        if session.get("done"):
            start = session.get("slide_start", 0)
            end = session.get("slide_end", 0)
            covered.update(range(start, end + 1))
    return all(i in covered for i in range(1, total_slides + 1))